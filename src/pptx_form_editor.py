from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import uuid
import json
import datetime
import base64
import os
from src.utils import logger

import re


class TextCleaner:
    """text border cleaner - Specially handledPPTXboundary characters in"""

    # Define boundary characters（Newlines, spaces, tabs, etc.）
    BOUNDARY_CHARS = r"[\n\r\t\f\v\u200b\u200c\u200d\u2060\s]*"

    @staticmethod
    def strip_boundary(text):
        """Strip border characters on both sides，return(core text, front boundary, posterior border)"""
        if not text:
            return "", "", ""

        # Match front boundary
        leading_match = re.match(f"^{TextCleaner.BOUNDARY_CHARS}", text)
        leading = leading_match.group() if leading_match else ""

        # Match back boundary
        trailing_match = re.search(f"{TextCleaner.BOUNDARY_CHARS}$", text)
        trailing = trailing_match.group() if trailing_match else ""

        # Extract core text
        start = len(leading)
        end = len(text) - len(trailing)
        core_text = text[start:end]

        return core_text, leading, trailing

    @staticmethod
    def restore_boundary(core_text, leading, trailing):
        """Restore boundary characters"""
        return leading + core_text + trailing


class PPTXFormEditor:
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path

    def extract_editable_content(self, assets_dir=None):
        """extractPPTXContent and identifies editable areas

        Args:
            assets_dir: Picture saving directory，If not provided, the image file will not be saved.
        """
        try:
            logger.debug(f"Parsing file: {self.pptx_path}")
            prs = Presentation(self.pptx_path)
            logger.debug(f"turn up {len(prs.slides)} slides")

            content_structure = {
                "slides": [],
                "session_id": str(uuid.uuid4()),
                "total_slides": len(prs.slides),
                "images": {},  # New：Store all picture information
            }

            for slide_idx, slide in enumerate(prs.slides):
                logger.debug(f"processing section {slide_idx + 1} slides...")
                slide_data = {
                    "slide_number": slide_idx + 1,
                    "layout": slide.slide_layout.name,
                    "shapes": [],
                }

                shape_count = 0
                for shape_idx, shape in enumerate(slide.shapes):

                    shape_info = self._extract_shape_info(
                        shape,
                        slide_idx,
                        shape_count,
                        parent_group=None,
                        assets_dir=assets_dir,
                    )
                    if shape_info:
                        slide_data["shapes"].append(shape_info)
                        shape_count += 1

                        # 🔧 Record picture information
                        if shape_info.get("has_image") and shape_info.get("image_ref"):
                            image_id = f"slide_{slide_idx}_image_{shape_idx}"
                            content_structure["images"][image_id] = {
                                "shape_id": shape_info["id"],
                                "filename": shape_info.get("image_ref"),
                                "info": shape_info.get("image_info", {}),
                            }

                        logger.debug(
                            f"  find shape {shape_idx}: {shape_info['type']} - '{shape_info.get('text', '')[:30]}...'"
                        )

                content_structure["slides"].append(slide_data)
                logger.debug(
                    f"No. {slide_idx + 1} slides found {len(slide_data['shapes'])} shapes"
                )

            # 🔧 debug：Check image references
            logger.debug("🔍 Check image citation information:")
            for slide_data in content_structure["slides"]:
                for shape in slide_data["shapes"]:
                    if shape.get("has_image"):
                        logger.debug(
                            f"  Picture shape {shape['id']}: image_ref={shape.get('image_ref')}, has_image_data={bool(shape.get('image_data'))}"
                        )
            return content_structure
        except Exception as e:
            logger.debug(f"Parse error: {e}")
            import traceback

            traceback.print_exc()
            return {"slides": [], "session_id": str(uuid.uuid4()), "total_slides": 0}

    def _extract_shape_info(
        self, shape, slide_idx, shape_idx, parent_group=None, assets_dir=None
    ):
        """Extract shape information"""
        try:
            shape_id = f"slide_{slide_idx}_shape_{shape_idx}"

            # Get basic shape information
            shape_name = getattr(shape, "name", f"shape_{shape_idx}")
            shape_type = self._get_shape_type_name(shape)

            logger.debug(f"    Work with shapes: {shape_name}, type: {shape_type}")

            # Get detailed location and size information
            shape_info = {
                "id": shape_id,
                "type": shape_type,
                "name": shape_name,
                "parent_group": parent_group,  # record parent combination
                "left": int(shape.left) if hasattr(shape, "left") else 100,
                "top": int(shape.top) if hasattr(shape, "top") else 100,
                "width": int(shape.width) if hasattr(shape, "width") else 200,
                "height": int(shape.height) if hasattr(shape, "height") else 50,
                "rotation": getattr(shape, "rotation", 0),
                "z_order": shape_idx,
            }
            # Extract text content
            text_obj = self._extract_text_content(shape)
            shape_info["text"] = (
                text_obj["paragraphs"] if "paragraphs" in text_obj else []
            )

            # 🔧 repair：Check if it is an image shape
            if self._is_picture_shape(shape):
                # Extract image information（incomingassets_dir）
                picture_info = self._extract_picture_info(
                    shape, slide_idx, shape_idx, assets_dir
                )
                if picture_info:
                    shape_info.update(picture_info)

            # Processing forms
            elif hasattr(shape, "has_table") and shape.has_table:
                table_info = self._extract_table_info(shape)
                if table_info:
                    shape_info.update(table_info)

            # Working with combined shapes
            elif shape_type == "group" and hasattr(shape, "shapes"):
                shape_info["is_group"] = True
                shape_info["child_shapes"] = []

                # Record the boundary information of the combination
                shape_info["group_bounds"] = {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }

                for child_idx, child_shape in enumerate(shape.shapes):
                    child_info = self._extract_shape_info(
                        child_shape,
                        slide_idx,
                        f"{shape_idx}_child_{child_idx}",
                        parent_group=shape_id,  # Pass parent combinationID
                    )
                    if child_info:
                        # Adjust the relative coordinates of a subshape
                        child_info["left"] = (
                            shape_info["left"] - 200 + (child_shape.left - shape.left)
                        )
                        child_info["top"] = (
                            shape_info["top"] - 200 + (child_shape.top - shape.top)
                        )
                        shape_info["child_shapes"].append(child_info)

            # Working with charts
            elif shape_type == "chart":
                chart_info = self._extract_chart_info(shape)
                if chart_info:
                    shape_info.update(chart_info)

            return shape_info

        except Exception as e:
            logger.debug(f"❌ Error in extracting shape information: {e}")
            import traceback

            traceback.print_exc()
            return None

    def _get_shape_type_name(self, shape):
        """Get shape type name"""
        try:
            raw_type = getattr(shape, "shape_type", "unknown")

            # Map numeric types to readable names
            type_mapping = {
                13: "picture",  # picture
                14: "placeholder",  # placeholder
                17: "textbox",  # text box
                19: "table",  # sheet
                1: "autoshape",  # automatic shape
                6: "group",  # combination
                9: "line",  # line
                3: "chart",  # chart
                5: "media",  # media
                8: "smartart",  # SmartArt
                18: "ole_object",  # OLEobject
            }

            if isinstance(raw_type, int):
                return type_mapping.get(raw_type, "unknown")
            elif isinstance(raw_type, str):
                # deal with "picture (13)" this situation
                return raw_type.split("(")[0].strip().lower()
            else:
                return "unknown"

        except Exception as e:
            logger.debug(f"Getting shape type error: {e}")
            return "unknown"

    def _is_picture_shape(self, shape):
        """Check whether it is an image shape"""
        try:
            # method1：examineshape_type
            if hasattr(shape, "shape_type"):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    return True

            # method2：Check if there isimageproperty
            if hasattr(shape, "image") and shape.image:
                return True

            # method3：Check if the name contains image keywords
            shape_name = getattr(shape, "name", "").lower()
            if any(
                keyword in shape_name
                for keyword in ["picture", "image", "photo", "img"]
            ):
                return True

            return False
        except:
            return False

    def _save_image_to_file(self, shape, slide_idx, shape_idx, assets_dir):
        """Save image to file，Return picture information"""
        try:
            if not hasattr(shape, "image") or not shape.image:
                return None

            # Get image binary data
            image_blob = shape.image.blob
            if not image_blob:
                return None

            # Get image format
            content_type = getattr(shape.image, "content_type", "image/png")
            ext_map = {
                "image/jpeg": ".jpg",
                "image/png": ".png",
                "image/gif": ".gif",
                "image/bmp": ".bmp",
            }
            extension = ext_map.get(content_type, ".png")

            # Generate unique file name
            import hashlib
            import uuid

            # Use content hashing + UUID ensure uniqueness
            content_hash = hashlib.md5(image_blob).hexdigest()[:8]
            unique_id = str(uuid.uuid4())[:8]
            filename = f"img_{slide_idx}_{shape_idx}{extension}"

            # Save picture file
            image_path = os.path.join(assets_dir, filename)
            with open(image_path, "wb") as f:
                f.write(image_blob)

            # Return picture information
            return {
                "filename": filename,
                "path": image_path,
                "size_bytes": len(image_blob),
                "width": (
                    shape.width.emu if hasattr(shape.width, "emu") else shape.width
                ),
                "height": (
                    shape.height.emu if hasattr(shape.height, "emu") else shape.height
                ),
                "content_type": content_type,
                "extension": extension,
            }

        except Exception as e:
            logger.debug(f"❌ Failed to save image: {e}")
            return None

    def _extract_picture_info(
        self, shape, slide_idx=None, shape_idx=None, assets_dir=None
    ):
        """Extract image information - Make sure to always return image_ref"""
        try:
            picture_info = {
                "has_image": True,
                "image_ref": None,  # This must be set
                "image_data": None,
                "image_info": {},
            }

            # 🔧 key：Generate base file name
            base_filename = f"img_{slide_idx}_{shape_idx}"

            # If a save directory is provided，then save as file
            if assets_dir and slide_idx is not None and shape_idx is not None:
                try:
                    saved_info = self._save_image_to_file(
                        shape, slide_idx, shape_idx, assets_dir
                    )
                    if saved_info:
                        picture_info["image_info"] = saved_info
                        picture_info["image_ref"] = saved_info[
                            "filename"
                        ]  # Use saved file name
                        logger.debug(
                            f"✅ save image: {saved_info['filename']} ({saved_info['size_bytes']} bytes)"
                        )
                    else:
                        # Save failed，Use base file name
                        picture_info["image_ref"] = f"{base_filename}.jpg"
                        logger.debug(
                            f"⚠️ Failed to save image，Use default references: {picture_info['image_ref']}"
                        )
                except Exception as e:
                    logger.debug(f"⚠️ Exception when saving picture file: {e}")
                    picture_info["image_ref"] = f"{base_filename}.jpg"
            else:
                # No save directory，Also set a reference
                picture_info["image_ref"] = f"{base_filename}.jpg"
                logger.debug(
                    f"ℹ️ No save directory，Use references: {picture_info['image_ref']}"
                )

            # Fallback tobase64coding（small picture），but keep the reference
            if hasattr(shape, "image") and shape.image:
                try:
                    image_bytes = shape.image.blob
                    if image_bytes and len(image_bytes) < 1000000:  # 1MBthe following
                        image_base64 = base64.b64encode(image_bytes).decode("utf-8")
                        content_type = getattr(shape.image, "content_type", "image/png")
                        picture_info["image_data"] = (
                            f"data:{content_type};base64,{image_base64}"
                        )

                        # Even if there isbase64data，Also keep file references
                        if not picture_info.get("image_ref"):
                            picture_info["image_ref"] = f"{base_filename}.jpg"
                except Exception as e:
                    logger.debug(f"⚠️ Base64Encoding failed: {e}")

            # 🔧 make sure image_ref Not empty
            if not picture_info["image_ref"]:
                picture_info["image_ref"] = f"{base_filename}.jpg"
                logger.debug(
                    f"🔧 make sure image_ref Not empty: {picture_info['image_ref']}"
                )

            logger.debug(
                f"📸 Image information extraction completed: ref={picture_info['image_ref']}, has_data={bool(picture_info['image_data'])}"
            )

            return picture_info

        except Exception as e:
            logger.debug(f"❌ Error extracting image information: {e}")
            # Return a reference even on error
            base_filename = (
                f"img_{slide_idx}_{shape_idx}" if slide_idx is not None else "unknown"
            )
            return {
                "has_image": True,
                "image_ref": f"{base_filename}.jpg",
                "image_error": str(e),
            }

    def _extract_table_info(self, shape):
        """Extract table information"""
        try:
            if hasattr(shape, "has_table") and shape.has_table:
                return {
                    "is_table": True,
                    "table_data": self._extract_table_content(shape.table),
                }
        except Exception as e:
            logger.debug(f"Error extracting form information: {e}")
        return None

    def _extract_chart_info(self, shape):
        """Extract chart information"""
        try:
            return {
                "is_chart": True,
                "chart_type": getattr(shape.chart, "chart_type", "unknown"),
                "chart_title": getattr(shape.chart, "chart_title", {})
                .get("text_frame", {})
                .get("text", ""),
            }
        except Exception as e:
            logger.debug(f"Error extracting chart information: {e}")
            return {"is_chart": True, "chart_error": str(e)}

    def _extract_text_content(self, shape):
        """Extract text content with full run-level structure preservation"""
        # Extract from text frame with run preservation
        if hasattr(shape, "text_frame"):
            paragraphs_data = []
            flat_text_parts = []
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                para_data = {
                    "paragraph_index": para_idx,
                    "alignment": (
                        str(paragraph.alignment) if paragraph.alignment else "left"
                    ),
                    "level": paragraph.level if hasattr(paragraph, "level") else 0,
                    "font": self._extract_paragraph_font(
                        paragraph
                    ),  # New：paragraph level font
                    "runs": [],
                }

                for run_idx, run in enumerate(paragraph.runs):
                    # Extract detailedrunformat information
                    core_text, leading, trailing = TextCleaner.strip_boundary(run.text)

                    run_data = {
                        "run_index": run_idx,
                        "text": core_text,  # 🎯 Only store core text
                        "boundary": {  # 🎯 Store boundary information
                            "leading": leading,
                            "trailing": trailing,
                        },
                        "editable": self._is_editable_text(core_text),
                        "format": self._extract_run_format(run),
                    }

                    para_data["runs"].append(run_data)

                paragraphs_data.append(para_data)

            # Store structured data on shape for later use
            shape._pptxpress_paragraphs_data = paragraphs_data
            # Return full hierarchy instead of flat text
            return {
                "type": "rich_text",
                "paragraphs": paragraphs_data,
                "shape_type": "text_frame",
            }

        # For tables, we also need to preserve run structure
        if hasattr(shape, "has_table") and shape.has_table:
            table_result = self._extract_table_content(shape.table)
            return {"type": "table", "content": table_result}

        return {"type": "empty", "text": ""}

    def _extract_paragraph_font(self, paragraph):
        """Extract paragraph-level font formatting"""
        font_data = {}

        try:
            # Paragraphs may also have font settings（as default）
            if hasattr(paragraph, "font"):
                font = paragraph.font
                font_data["bold"] = getattr(font, "bold", None)
                font_data["italic"] = getattr(font, "italic", None)
                font_data["size"] = getattr(font, "size", None)
                font_data["name"] = getattr(font, "name", None)
                font_data["color"] = self._extract_color(getattr(font, "color", None))
        except:
            pass

        return font_data

    def _extract_color(self, color_obj):
        """Extract color information"""
        if color_obj is None:
            return None

        color_info = {"type": str(type(color_obj))}

        try:
            # RGBcolor
            if hasattr(color_obj, "rgb"):
                color_info["rgb"] = color_obj.rgb
                # Convert to hexadecimal
                if color_obj.rgb:
                    color_info["hex"] = f"#{color_obj.rgb:06X}"

            # theme color
            if hasattr(color_obj, "theme_color"):
                color_info["theme_color"] = str(color_obj.theme_color)

            # Brightness adjustment
            if hasattr(color_obj, "brightness"):
                color_info["brightness"] = color_obj.brightness

            # color type
            if hasattr(color_obj, "type"):
                color_info["color_type"] = str(color_obj.type)

        except Exception as e:
            color_info["error"] = str(e)

        return color_info

    def _extract_table_content(self, table):
        """Extract table content"""
        try:
            table_data = []
            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell_idx, cell in enumerate(row.cells):
                    cell_text = ""
                    if cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                cell_text += run.text
                    row_data.append(
                        {"text": cell_text.strip(), "row": row_idx, "col": cell_idx}
                    )
                table_data.append(row_data)
            return json.dumps(table_data, ensure_ascii=False)
        except Exception as e:
            logger.debug(f"Extract form error: {e}")
            return "[]"

    def _is_editable_text(self, text):
        """Determine whether text is editable"""
        if not text.strip():
            return True

        editable_indicators = [
            "Click to enter",
            "Please enter",
            "Fill in the content",
            "Add title",
            "Add text",
            "subtitle",
            "Click to add",
            "Enter text",
            "Add content",
            "Title",
            "Subtitle",
            "Click to edit",
            "Type here",
            "Your text here",
            "Text here",
        ]

        # Short text or text containing indicators can be edited
        text_lower = text.lower()
        return len(text.strip()) < 10 or any(
            indicator.lower() in text_lower for indicator in editable_indicators
        )

    def generate_editable_html(self, content_structure):
        """Generate editableHTMLinterface"""
        logger.debug(f"🔍 generateHTML: {len(content_structure['slides'])}slides")

        # Decide which content template to display
        if not content_structure["slides"] or all(
            len(slide["shapes"]) == 0 for slide in content_structure["slides"]
        ):
            content_html = template_loader.load_template("empty_content.html")
            javascript_data = 'console.log("No content to edit");'
        else:
            content_html = template_loader.render_template(
                "editor_content.html", total_slides=content_structure["total_slides"]
            )
            javascript_data = self._generate_javascript_data(content_structure)

        # Render base template
        html_content = template_loader.render_template(
            "base.html", content=content_html, javascript_data=javascript_data
        )

        return html_content

    def _generate_javascript_data(self, content_structure):
        """generateJavaScriptinitialization data"""
        try:
            # compactJSONFormat
            slides_data_json = json.dumps(
                content_structure["slides"], ensure_ascii=False, separators=(",", ":")
            )

            js_code = f"""
            // Initialize editor data
            window.editorData = {{
                slidesData: {slides_data_json},
                sessionId: "{content_structure['session_id']}",
                totalSlides: {content_structure['total_slides']}
            }};
            console.log("✅ Data setting completed");
            """

            return js_code

        except Exception as e:
            logger.debug(f"❌ generateJavaScriptData error: {e}")
            return 'console.error("Data generation error");'

    def apply_changes_to_pptx(self, presentation, changes, content_structure):
        """Apply changes toPPTXdocument"""
        if not changes:
            logger.debug("⚠️ No modifications need to be applied")
            return

        logger.debug(f"🔄 Start applying {len(changes)} Modify everywhere...")

        applied_count = 0
        for shape_id, new_text in changes.items():
            if self._apply_single_change(
                presentation, shape_id, new_text, content_structure
            ):
                applied_count += 1

        logger.debug(
            f"✅ Successfully applied {applied_count}/{len(changes)} Modify everywhere"
        )

    def _apply_single_change(
        self, presentation, shape_id, change_data, content_structure
    ):
        """Apply single modification"""
        try:
            # 🔧 repair：Processing dictionary formatchange_data
            if isinstance(change_data, dict):
                new_text = change_data.get("text", "")
                if not new_text and "txt" in change_data:
                    new_text = change_data.get("txt", "")
            else:
                new_text = str(change_data)

            # 🔧 repair：Skip empty or placeholder text
            if not new_text.strip():
                logger.debug(f"⏭️ Skip empty text: {shape_id}")
                return False

            # Filter placeholder text
            placeholder_patterns = [
                "[placeholder]",
                "◇ shape",
                "debug",
                "No content",
                "placeholder",
                "shape",
            ]
            if any(pattern in new_text for pattern in placeholder_patterns):
                logger.debug(
                    f"⏭️ Skip placeholder text: {shape_id} -> '{new_text[:50]}...'"
                )
                return False

            # Analyze shapesID
            parts = shape_id.split("_")
            if len(parts) >= 4 and parts[0] == "slide" and parts[2] == "shape":
                slide_idx = int(parts[1])
                shape_idx = int(parts[3])

                # Verify index
                if 0 <= slide_idx < len(presentation.slides):
                    slide = presentation.slides[slide_idx]

                    if 0 <= shape_idx < len(slide.shapes):
                        shape = slide.shapes[shape_idx]

                        # 🔧 repair：Check if the text shape is editable
                        if not self._is_text_shape(shape):
                            shape_type_name = self._get_shape_type_name(shape)
                            logger.debug(
                                f"⏭️ Skip non-text shapes: {shape_id}, type: {shape_type_name}"
                            )
                            return False

                        # Apply changes
                        success = self._update_shape_content(shape, new_text)
                        if success:
                            logger.debug(
                                f"✅ Apply changes: {shape_id} -> '{new_text[:50]}...'"
                            )
                        return success
                    else:
                        logger.debug(f"⚠️ Invalid shape index: {shape_idx}")
                else:
                    logger.debug(f"⚠️ Invalid slide index: {slide_idx}")
            else:
                logger.debug(f"⚠️ Unable to parse shapeID: {shape_id}")

        except Exception as e:
            logger.debug(f"❌ Applying changes failed {shape_id}: {e}")
            import traceback

            traceback.print_exc()

        return False

    def _is_text_shape(self, shape):
        """Check if the text shape is editable"""
        try:
            shape_type = getattr(shape, "shape_type", None)

            # Text-related shape types
            text_shape_types = [
                14,  # placeholder PLACEHOLDER
                17,  # text box TEXT_BOX
                1,  # automatic shape（may contain text）AUTO_SHAPE
            ]

            # method1：Check shape type
            if shape_type in text_shape_types:
                return True

            # method2：Check if there is a text frame
            if hasattr(shape, "text_frame") and shape.text_frame:
                # Check if text frame has content
                try:
                    if shape.text_frame.text and shape.text_frame.text.strip():
                        return True
                except:
                    pass

            # method3：Check if there istextproperty
            if hasattr(shape, "text") and shape.text is not None:
                return True

            # Non-text shape types
            non_text_shape_types = [
                13,  # picture PICTURE
                6,  # combination GROUP
                9,  # line LINE
                3,  # chart CHART
                5,  # media MEDIA
                8,  # SmartArt SMART_ART
                18,  # OLEobject OLE_OBJECT
                19,  # sheet TABLE
            ]

            if shape_type in non_text_shape_types:
                return False

            # Not modified by default
            return False

        except Exception as e:
            logger.debug(f"Checking shape type failed: {e}")
            return False

    def _update_shape_content(self, shape, new_text):
        """Update shape content"""
        try:
            # text frame shape（Text boxes, placeholders, etc.）
            if hasattr(shape, "text_frame") and shape.text_frame:
                self._update_text_frame(shape.text_frame, new_text)
                return True

            # table shape
            elif hasattr(shape, "has_table") and shape.has_table:
                # Table modification requires special handling，Simplified processing here
                logger.debug(
                    f"📊 Table modification requires special handling: {shape.name}"
                )
                return False

            # Auto shapes may contain text
            elif hasattr(shape, "text"):
                shape.text = new_text
                return True

            else:
                logger.debug(f"⚠️ Modified shape type not supported: {shape.shape_type}")
                return False

        except Exception as e:
            logger.debug(f"❌ Failed to update shape content: {e}")
            return False

    def _update_text_frame(self, text_frame, new_text):
        """Update text frame content"""
        # Clear all existing text
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = ""

        # If there is a paragraph，Add new text using first paragraph
        if text_frame.paragraphs:
            first_paragraph = text_frame.paragraphs[0]
            if first_paragraph.runs:
                first_paragraph.runs[0].text = new_text
            else:
                first_paragraph.add_run().text = new_text
        else:
            # if there is no paragraph，create new
            text_frame.paragraphs[0].add_run().text = new_text

    def extract_for_project_init(self, project_dir: str = None) -> dict:
        """
        Extract data specifically for project initialization，Returns all required formats.

        Args:
            project_dir: Project directory，If provided save the image toassetsTable of contents

        Returns:
            dict: Contains data in various formats
        """
        try:
            logger.debug(f"🔄 Extract project initialization data: {self.pptx_path}")

            # SureassetsTable of contents
            assets_dir = None
            if project_dir:
                assets_dir = os.path.join(project_dir, "assets", "images")
                os.makedirs(assets_dir, exist_ok=True)
                logger.debug(f"📁 The picture will be saved to: {assets_dir}")

            # 1. Extract original content（incomingassets_dir）
            raw_content = self.extract_editable_content(assets_dir=assets_dir)

            if not raw_content or "slides" not in raw_content:
                logger.debug("⚠️ No slide content was extracted")
                return {
                    "raw_content": {"slides": []},
                    "editor_format": {"slides": []},
                    "slide_files_data": [],
                    "slides_info": [],
                    "images_info": {},
                }

            # 2. Generate data in various formats
            editor_format = self._to_editor_format(raw_content)
            slide_files_data = self._to_slide_files_data(raw_content)
            slides_info = self._to_slides_info(slide_files_data)
            logger.debug(f"✅ Extraction completed:")
            logger.debug(
                f"   - original content: {len(raw_content.get('slides', []))} slides"
            )
            logger.debug(
                f"   - Number of pictures: {len(raw_content.get('images', {}))} open"
            )
            logger.debug(
                f"   - Editor format: {len(editor_format.get('slides', []))} slides"
            )
            logger.debug(f"   - slidefile data: {len(slide_files_data)} indivual")
            logger.debug(f"   - slidesinformation: {len(slides_info)} indivual")

            return {
                "raw_content": raw_content,
                "editor_format": editor_format,
                "slide_files_data": slide_files_data,
                "slides_info": slides_info,
                "images_info": raw_content.get("images", {}),  # New
            }

        except Exception as e:
            logger.debug(f"❌ Failed to extract project initialization data: {e}")
            import traceback

            traceback.print_exc()
            return {
                "raw_content": {"slides": []},
                "editor_format": {"slides": []},
                "slide_files_data": [],
                "slides_info": [],
                "images_info": {},
            }

    def _to_editor_format(self, raw_content: dict) -> dict:
        """
        Convert raw content into the format expected by the front-end editor.

        Args:
            raw_content: extract_editable_content() return value

        Returns:
            dict: Front-end editor format
        """
        editor_slides = []

        for slide_data in raw_content.get("slides", []):
            slide_obj = {
                "slide_number": slide_data.get("slide_number", 0),
                "shapes": [],
            }

            for shape in slide_data.get("shapes", []):
                shape_obj = {
                    "id": shape.get("id", ""),
                    "type": shape.get("type", "placeholder"),
                    "name": shape.get("name", ""),
                    "text": shape.get("text", ""),
                    "left": shape.get("left", 0),
                    "top": shape.get("top", 0),
                    "width": shape.get("width", 100),
                    "height": shape.get("height", 50),
                    "editable": shape.get("editable", True),
                }

                if shape.get("is_table"):
                    shape_obj["is_table"] = True
                    shape_obj["table_data"] = shape.get("table_data", "[]")
                # 🔧 critical fix：Pass image reference information
                if shape.get("has_image"):
                    shape_obj["has_image"] = True
                    shape_obj["image_data"] = shape.get("image_data")

                    # 🔧 transfer image_ref
                    if shape.get("image_ref"):
                        shape_obj["image_ref"] = shape.get("image_ref")
                    elif shape.get("image_info") and shape.get("image_info").get(
                        "filename"
                    ):
                        # from image_info Get the file name in
                        shape_obj["image_ref"] = shape.get("image_info").get("filename")

                slide_obj["shapes"].append(shape_obj)

            editor_slides.append(slide_obj)

        return {"slides": editor_slides}

    def _to_slide_files_data(self, raw_content: dict) -> list:
        """
        Convert original content toslideFile data format.

        Args:
            raw_content: extract_editable_content() return value

        Returns:
            list: slideFile data list
        """
        slide_files_data = []

        for slide_idx, slide_data in enumerate(raw_content.get("slides", [])):
            slide_id = f"slide_{slide_idx:03d}"
            shapes_data = {}

            # Process each shape
            for shape in slide_data.get("shapes", []):
                shape_id = shape.get("id", "")
                shape_type_abbr = self._get_shape_type_abbr(shape)

                shape_entry = {
                    "t": shape_type_abbr,
                    "mod": datetime.datetime.now().isoformat(),
                }

                # If it is a picture，Record picture information
                if shape.get("is_table"):
                    shape_entry["t"] = "table"
                    shape_entry["is_table"] = True
                    shape_entry["table_data"] = shape.get(
                        "table_data", "[]"
                    )  # 🎯 JSON str
                elif shape.get("has_image"):
                    shape_entry["type"] = "image"
                    if shape.get("image_ref"):
                        shape_entry["image_ref"] = shape.get("image_ref")
                    if shape.get("image_info"):
                        shape_entry["image_info"] = shape.get("image_info")
                else:
                    # text shape
                    shape_entry["txt"] = shape.get("text", "")

                shapes_data[shape_id] = shape_entry

            # Extract title
            title = self._extract_slide_title(shapes_data)

            slide_files_data.append(
                {
                    "slide_id": slide_id,
                    "slide_number": slide_idx + 1,
                    "shapes_data": shapes_data,
                    "title": title,
                }
            )

        return slide_files_data

    def _to_slides_info(self, slide_files_data: list) -> list:
        """
        fromslideFile data is extracted forproject.yamlinformation.

        Args:
            slide_files_data: _to_slide_files_data() return value

        Returns:
            list: used forproject.yamlofslideInformation list
        """
        slides_info = []

        for slide_data in slide_files_data:
            slides_info.append(
                {
                    "id": slide_data["slide_id"],
                    "file": f"slides/{slide_data['slide_id']}.json",
                    "title": slide_data["title"][:50],  # Limit title length
                }
            )

        return slides_info

    def _get_shape_type_abbr(self, shape: dict) -> str:
        """
        Get the abbreviation based on the shape type.

        Args:
            shape: shape data dictionary

        Returns:
            str: Shape type abbreviation
        """
        full_type = shape.get("type", "").lower()

        # type mapping
        if "title" in full_type:
            return "ttl"
        elif "subtitle" in full_type:
            return "sub"
        elif "textbox" in full_type or "placeholder" in full_type:
            return "txt"
        else:
            return "txt"  # default

    def _extract_slide_title(self, shapes_data: dict) -> str:
        """
        Extract slide titles from shape data.

        Args:
            shapes_data: shape data dictionary {shape_id: shape_info}

        Returns:
            str: slide title
        """
        if not shapes_data:
            return ""

        # First find the shape of the title type
        for shape_id, shape_info in shapes_data.items():
            if shape_info.get("t") == "ttl":
                title = shape_info.get("txt", "")
                if title.strip():
                    return title

        # if title not found，Get the contents of the first text box
        first_shape = next(iter(shapes_data.values()), {})
        return first_shape.get("txt", "")

    def save_slide_files(self, slide_files_data: list, output_dir: str) -> bool:
        """
        saveslidefile to the specified directory.

        Args:
            slide_files_data: _to_slide_files_data() return value
            output_dir: Output directory（Project directory）

        Returns:
            bool: Is the save successful?
        """
        try:
            slides_dir = os.path.join(output_dir, "slides")
            os.makedirs(slides_dir, exist_ok=True)

            for slide_data in slide_files_data:
                slide_id = slide_data["slide_id"]
                slide_json = {
                    "id": slide_id,
                    "slide_number": slide_data["slide_number"],
                    "shapes": slide_data["shapes_data"],
                }

                for shape_id, shape_info in slide_data["shapes_data"].items():
                    if shape_info.get("t") == "table" and "table_data" in shape_info:
                        # Change table_data to str
                        if isinstance(shape_info["table_data"], list):
                            shape_info["table_data"] = json.dumps(
                                shape_info["table_data"], ensure_ascii=False
                            )

                # save asJSONdocument
                slide_path = os.path.join(slides_dir, f"{slide_id}.json")
                with open(slide_path, "w", encoding="utf-8") as f:
                    json.dump(slide_json, f, separators=(",", ":"), ensure_ascii=False)

                logger.debug(f"💾 saveslidedocument: {slide_path}")

            logger.debug(
                f"✅ Save completed: {len(slide_files_data)} indivualslidedocument"
            )
            return True

        except Exception as e:
            logger.debug(f"❌ saveslideFile failed: {e}")
            import traceback

            traceback.print_exc()
            return False

    def apply_image_to_pptx(self, presentation, shape_id, image_path):
        """Apply image toPPTXspecified shape"""
        try:
            # Analyze shapesID
            parts = shape_id.split("_")
            if len(parts) >= 4 and parts[0] == "slide" and parts[2] == "shape":
                slide_idx = int(parts[1])
                shape_idx = int(parts[3])

                if 0 <= slide_idx < len(presentation.slides):
                    slide = presentation.slides[slide_idx]

                    if 0 <= shape_idx < len(slide.shapes):
                        shape = slide.shapes[shape_idx]

                        # Check if it is an image shape
                        if self._is_picture_shape(shape):
                            # Replace picture
                            return self._replace_shape_image(shape, image_path)
                        else:
                            logger.debug(f"⚠️ Shape is not an image type: {shape_id}")
                    else:
                        logger.debug(f"⚠️ Invalid shape index: {shape_idx}")
                else:
                    logger.debug(f"⚠️ Invalid slide index: {slide_idx}")

            return False

        except Exception as e:
            logger.debug(f"❌ Applying image failed: {e}")
            return False

    def _replace_shape_image(self, shape, image_path):
        """Replace image with shape - use the correctpython-pptx API"""
        try:
            if not os.path.exists(image_path):
                logger.debug(f"❌ Image file does not exist: {image_path}")
                return False

            logger.debug(
                f"🖼️ Replace picture: {shape.name} <- {os.path.basename(image_path)}"
            )

            # method1：using pictures part property
            try:
                # Get picturesrelation id
                rId = shape._pic.blipFill.blip.rEmbed

                # passslide.partGet relevant parts
                slide = shape._parent
                image_part = slide.part.related_part(
                    rId
                )  # 🔧 use related_part，no related_parts

                # Read new image data
                with open(image_path, "rb") as f:
                    new_image_data = f.read()

                # Replace image data
                image_part._blob = new_image_data

                logger.debug(
                    f"✅ Image replacement successful（method1）: {shape.name}"
                )
                return True

            except AttributeError as e:
                logger.debug(f"⚠️ method1fail: {e}")

                # method2：A simpler and more direct method
                try:
                    # Get the picture part directly
                    pic = shape._pic
                    rId = pic.blipFill.blip.rEmbed

                    # Get relationship
                    rel = pic.part.rels[rId]

                    # Read new pictures
                    with open(image_path, "rb") as f:
                        new_image_data = f.read()

                    # direct replacement
                    rel._target._blob = new_image_data

                    logger.debug(
                        f"✅ Image replacement successful（method2）: {shape.name}"
                    )
                    return True

                except Exception as e2:
                    logger.debug(f"⚠️ method2fail: {e2}")

                    # method3：Use image replacement method（if there is）
                    try:
                        return self._replace_image_direct(shape, image_path)
                    except Exception as e3:
                        logger.debug(f"❌ method3fail: {e3}")
                        return False

        except Exception as e:
            logger.debug(f"❌ Failed to replace image: {e}")
            import traceback

            traceback.print_exc()
            return False

    def _extract_run_format(self, run):
        """Extract format information - SIMPLIFIED & SERIALIZABLE"""
        font = run.font

        # Base format info
        format_info = {
            "bold": bool(font.bold) if font.bold is not None else False,
            "italic": bool(font.italic) if font.italic is not None else False,
            "underline": bool(font.underline) if hasattr(font, "underline") else False,
            "size": float(font.size.pt) if font.size else None,
            "font": str(font.name) if font.name else None,
            "color": self._get_color_info(font.color),
        }

        # Handle language safely
        if hasattr(font, "language_id") and font.language_id is not None:
            try:
                # Convert to simple int or string
                format_info["language"] = int(font.language_id)
            except:
                format_info["language"] = str(font.language_id)

        # Convert all values to basic Python types
        return self._make_serializable(format_info)

    def _make_serializable(self, obj):
        """Recursively convert objects to YAML-serializable types"""
        if obj is None:
            return None
        elif isinstance(obj, (str, int, float, bool)):
            return obj
        elif isinstance(obj, (list, tuple)):
            return [self._make_serializable(item) for item in obj]
        elif isinstance(obj, dict):
            return {
                str(key): self._make_serializable(value) for key, value in obj.items()
            }
        else:
            # Convert any other type to string
            return str(obj)

    def _get_color_info(self, color):
        """Simplified color extraction - always returns RGB hex or None"""
        if color is None:
            return None

        try:
            # Method 1: Direct RGB access
            if hasattr(color, "rgb") and color.rgb:
                rgb = color.rgb
                if isinstance(rgb, tuple) and len(rgb) >= 3:
                    r, g, b = rgb[0], rgb[1], rgb[2]
                    # Ensure values are 0-255
                    r = max(0, min(255, int(r)))
                    g = max(0, min(255, int(g)))
                    b = max(0, min(255, int(b)))
                    return f"#{r:02x}{g:02x}{b:02x}"

            # Method 2: Check type
            if hasattr(color, "type"):
                if color.type == 1:  # RGB
                    rgb = getattr(color, "rgb", None)
                    if rgb and isinstance(rgb, tuple) and len(rgb) >= 3:
                        r, g, b = rgb[0], rgb[1], rgb[2]
                        return f"#{r:02x}{g:02x}{b:02x}"

                elif color.type == 2:  # Theme color
                    # 🎯 Convert theme colors to default RGB
                    # You can customize this mapping based on your template
                    theme_color = getattr(color, "theme_color", None)
                    if theme_color:
                        theme_name = str(theme_color)
                        # Map common theme colors to RGB
                        theme_to_rgb = {
                            "ACCENT_1": "#4F81BD",  # Blue
                            "ACCENT_2": "#C0504D",  # Red
                            "ACCENT_3": "#9BBB59",  # Green
                            "ACCENT_4": "#8064A2",  # Purple
                            "ACCENT_5": "#4BACC6",  # Light Blue
                            "ACCENT_6": "#F79646",  # Orange
                            "DARK_1": "#000000",  # Black
                            "DARK_2": "#1F497D",  # Dark Blue
                            "LIGHT_1": "#FFFFFF",  # White
                            "LIGHT_2": "#EEECE1",  # Light Gray
                        }

                        for key, rgb_value in theme_to_rgb.items():
                            if key in theme_name:
                                return rgb_value

                    # Default for theme colors
                    return "#000000"  # Black

            # Method 3: Try string parsing
            color_str = str(color)
            if "RGB(" in color_str:
                import re

                match = re.search(r"RGB\((\d+),\s*(\d+),\s*(\d+)\)", color_str)
                if match:
                    r, g, b = (
                        int(match.group(1)),
                        int(match.group(2)),
                        int(match.group(3)),
                    )
                    return f"#{r:02x}{g:02x}{b:02x}"

            # If all else fails, return black or None
            return "#000000"

        except Exception as e:
            logger.debug(f"Color extraction error: {e}")
            return "#000000"  # Black fallback

    def apply_structured_changes_to_pptx(self, presentation, structured_changes):
        """Apply structural changes toPPTX（one by oneRunrenew，protect borders）"""
        updated_shapes = 0
        updated_runs = 0

        for shape_id, structured_data in structured_changes.items():
            try:
                # Skip if this is a table (tables are handled separately)
                if shape_id.startswith("slide_") and "_shape_" in shape_id:
                    parts = shape_id.split("_")
                    if len(parts) >= 4:
                        slide_idx = int(parts[1])
                        shape_idx = int(parts[3])

                        if 0 <= slide_idx < len(presentation.slides):
                            slide = presentation.slides[slide_idx]
                            if 0 <= shape_idx < len(slide.shapes):
                                shape = slide.shapes[shape_idx]
                                if hasattr(shape, "has_table") and shape.has_table:
                                    logger.debug(
                                        f"⏭️ Skipping table shape in text update: {shape_id}"
                                    )
                                    continue

                # Analyze shapesID
                parts = shape_id.split("_")
                if not (
                    len(parts) >= 4 and parts[0] == "slide" and parts[2] == "shape"
                ):
                    continue

                slide_idx = int(parts[1])
                shape_idx = int(parts[3])

                # Get the corresponding shape
                if 0 <= slide_idx < len(presentation.slides):
                    slide = presentation.slides[slide_idx]
                    if 0 <= shape_idx < len(slide.shapes):
                        shape = slide.shapes[shape_idx]

                        # Update text
                        runs_updated = self._apply_text_with_boundary(
                            shape, structured_data
                        )
                        if runs_updated > 0:
                            updated_shapes += 1
                            updated_runs += runs_updated

            except Exception as e:
                logger.debug(f"❌ deal with {shape_id} fail: {e}")

        logger.debug(
            f"📊 Text update completed: {updated_shapes} shapes, {updated_runs} runs"
        )
        return updated_shapes

    def _apply_text_with_boundary(self, shape, structured_data):
        """one by oneRunUpdate text，protect borders"""
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            return False

        text_frame = shape.text_frame
        updated_runs = 0

        # Process by paragraph
        for para_idx, para_data in enumerate(structured_data):
            if para_idx >= len(text_frame.paragraphs):
                break

            paragraph = text_frame.paragraphs[para_idx]
            runs_data = para_data.get("runs", [])

            # one by oneRunrenew
            for run_idx, run_data in enumerate(runs_data):
                if run_idx >= len(paragraph.runs):
                    break

                run = paragraph.runs[run_idx]

                # 🎯 restore boundaries
                core_text = run_data.get("text", "")
                boundary = run_data.get("boundary", {})
                leading = boundary.get("leading", "")
                trailing = boundary.get("trailing", "")

                # Splice complete text
                full_text = leading + core_text + trailing
                run.text = full_text
                updated_runs += 1

        logger.debug(f"✅ Update completed: {updated_runs} runs")
        return updated_runs > 0

    def apply_table_changes_to_pptx(self, presentation, table_changes):
        """Apply table changes toPPTX（only update text，protect format）"""
        updated_tables = 0
        updated_cells = 0

        for shape_id, table_data in table_changes.items():
            try:
                # Parse shape ID
                parts = shape_id.split("_")
                if not (
                    len(parts) >= 4 and parts[0] == "slide" and parts[2] == "shape"
                ):
                    logger.debug(f"⚠️ Invalid shape ID format: {shape_id}")
                    continue

                slide_idx = int(parts[1])
                shape_idx = int(parts[3])

                # Get the corresponding shape
                if 0 <= slide_idx < len(presentation.slides):
                    slide = presentation.slides[slide_idx]
                    if 0 <= shape_idx < len(slide.shapes):
                        shape = slide.shapes[shape_idx]

                        # Check if it's a table
                        if shape.has_table:
                            # Update table cells
                            cells_updated = self._apply_table_cell_changes(
                                shape, table_data
                            )
                            if cells_updated > 0:
                                updated_tables += 1
                                updated_cells += cells_updated
                                logger.debug(
                                    f"✅ Updated table {shape_id}: {cells_updated} cells"
                                )
                        else:
                            logger.debug(f"⚠️ Shape {shape_id} is not a table")

            except Exception as e:
                logger.debug(f"❌ Failed to update table {shape_id}: {e}")

        logger.debug(
            f"📊 Table update completed: {updated_tables} tables, {updated_cells} cells"
        )
        return updated_tables

    def _apply_table_cell_changes(self, shape, table_data):
        """Update table cell text while preserving formatting"""
        if not hasattr(shape, "table") or not shape.table:
            return 0

        table = shape.table
        updated_cells = 0

        try:
            # Ensure table_data is a list of rows
            if not isinstance(table_data, list):
                logger.debug(f"⚠️ Table data is not a list: {type(table_data)}")
                return 0

            # Update each cell
            for row_idx, row_data in enumerate(table_data):
                if row_idx >= len(table.rows):
                    logger.debug(f"⚠️ Row index {row_idx} out of bounds")
                    break

                row = table.rows[row_idx]

                if isinstance(row_data, list):
                    for col_idx, cell_data in enumerate(row_data):
                        if col_idx >= len(row.cells):
                            logger.debug(
                                f"⚠️ Column index {col_idx} out of bounds in row {row_idx}"
                            )
                            break

                        cell = row.cells[col_idx]

                        if isinstance(cell_data, dict) and "text" in cell_data:
                            new_text = cell_data.get("text", "")
                            old_text = cell.text if hasattr(cell, "text") else ""

                            # Only update if text has changed
                            if new_text != old_text:
                                # Preserve formatting by updating text frame runs
                                if hasattr(cell, "text_frame") and cell.text_frame:
                                    self._update_table_cell_text(cell, new_text)
                                    updated_cells += 1
                                else:
                                    # Fallback: direct text update
                                    cell.text = new_text
                                    updated_cells += 1

        except Exception as e:
            logger.debug(f"❌ Error updating table cells: {e}")

        return updated_cells

    def _update_table_cell_text(self, cell, new_text):
        """Update table cell text while preserving run formatting"""
        if not hasattr(cell, "text_frame") or not cell.text_frame:
            cell.text = new_text
            return

        text_frame = cell.text_frame
        new_text = str(new_text) if new_text is not None else ""

        try:
            # Check if we need to preserve formatting
            if (
                len(text_frame.paragraphs) == 1
                and len(text_frame.paragraphs[0].runs) == 1
            ):
                # Single run, simple update
                run = text_frame.paragraphs[0].runs[0]
                run.text = new_text
            else:
                # Multiple runs - need to be careful
                # Strategy 1: Try to preserve first run's formatting for all text
                if text_frame.paragraphs:
                    first_para = text_frame.paragraphs[0]
                    if first_para.runs:
                        # Clear existing runs
                        for para in text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ""

                        # Update first run with new text
                        first_para.runs[0].text = new_text
                    else:
                        # No runs, use direct text
                        text_frame.text = new_text
                else:
                    # No paragraphs, use direct text
                    text_frame.text = new_text

        except Exception as e:
            logger.debug(f"⚠️ Error updating cell text with formatting: {e}")
            # Fallback to direct text update
            cell.text = new_text
